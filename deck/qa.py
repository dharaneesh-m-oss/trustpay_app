"""Geometry and encoding checks on the built deck.

Renders are the real visual QA, but these catch the two defects that are
objective rather than aesthetic: a shape that runs past the slide edge, and text
whose estimated width exceeds the box holding it.

Widths are estimated from average character advance for the font size, which is
approximate — it is deliberately conservative so it flags candidates for a human
look rather than pretending to be a layout engine.
"""

from __future__ import annotations

import sys
import zipfile

from pptx import Presentation
from pptx.util import Emu

DECK = sys.argv[1] if len(sys.argv) > 1 else "TrustPay.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.5  # minimum comfortable margin, inches

# Average advance width as a fraction of point size, for a humanist sans at
# mixed case. Bold and letter-spacing widen it.
AVG_ADVANCE = 0.50
BOLD_FACTOR = 1.055


def inches(v) -> float:
    return Emu(v).inches if v is not None else 0.0


def estimate_width(text: str, size_pt: float, bold: bool, spacing_pt: float) -> float:
    per_char = size_pt * AVG_ADVANCE * (BOLD_FACTOR if bold else 1.0) + spacing_pt
    return len(text) * per_char / 72.0


problems: list[str] = []
notes: list[str] = []

prs = Presentation(DECK)
for index, slide in enumerate(prs.slides, start=1):
    for shape in slide.shapes:
        x, y = inches(shape.left), inches(shape.top)
        w, h = inches(shape.width), inches(shape.height)

        if x < -0.01 or y < -0.01 or x + w > SLIDE_W + 0.01 or y + h > SLIDE_H + 0.01:
            problems.append(
                f"slide {index}: shape runs off the canvas "
                f"(x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f})"
            )
        elif x < MARGIN - 0.01 and w > 0.2:
            notes.append(f"slide {index}: starts at x={x:.2f}, inside the {MARGIN}\" margin")

        if not shape.has_text_frame:
            continue

        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs)
            if not line.strip():
                continue
            run = para.runs[0]
            size = run.font.size.pt if run.font.size else 18.0
            bold = bool(run.font.bold)
            est = estimate_width(line, size, bold, 0.0)
            # Only flag single-line boxes: a tall box is expected to wrap.
            lines_that_fit = max(1, int(h / (size * 1.25 / 72.0)))
            if est > w * lines_that_fit * 1.02:
                problems.append(
                    f"slide {index}: text may overflow its box "
                    f"(w={w:.2f}\" est={est:.2f}\" size={size:.0f}pt) :: {line[:64]!r}"
                )

# Encoding: the XML should carry real typographic characters, not mojibake.
with zipfile.ZipFile(DECK) as zf:
    blob = b"".join(
        zf.read(n) for n in zf.namelist() if n.startswith("ppt/slides/slide")
    )
text = blob.decode("utf-8")
for bad in ("�", "â€™", "â€œ", "Â "):
    if bad in text:
        problems.append(f"encoding: found {bad!r} in slide XML")

curly = text.count("’")
print(f"typographic apostrophes in XML: {curly}")
print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")

print("\n--- problems ---")
print("\n".join(problems) if problems else "  none")
print("\n--- notes ---")
print("\n".join(notes) if notes else "  none")
